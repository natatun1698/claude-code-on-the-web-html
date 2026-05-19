#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
AIエージェント・ストラテジスト 認定試験対策 プレゼンテーション生成スクリプト
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Color definitions
NAVY = RGBColor(0x1A, 0x23, 0x7E)       # Dark navy blue
GOLD = RGBColor(0xF5, 0x7F, 0x17)       # Gold/amber
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
DARK_GRAY = RGBColor(0x37, 0x47, 0x4F)
LIGHT_NAVY = RGBColor(0x3F, 0x51, 0xB5)  # Lighter navy for sub-headers
LIGHT_GOLD = RGBColor(0xFF, 0xB3, 0x00)  # Lighter gold

FONT_NAME = "Meiryo"

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

def add_background(slide, color=WHITE):
    """Add solid background to slide."""
    from pptx.util import Emu
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    """Add a rectangle shape."""
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()

    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()

    return shape

def add_text_box(slide, text, left, top, width, height,
                 font_size=18, font_color=DARK_GRAY, font_name=FONT_NAME,
                 bold=False, align=PP_ALIGN.LEFT, word_wrap=True):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.name = font_name
    run.font.bold = bold
    return txBox

def create_title_slide(prs):
    """Create the title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Full navy background
    add_background(slide, NAVY)

    # Gold accent bar at top
    add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.15), GOLD)

    # Gold accent bar at bottom
    add_rect(slide, Inches(0), Inches(7.35), Inches(13.33), Inches(0.15), GOLD)

    # Decorative left bar
    add_rect(slide, Inches(0.5), Inches(1.5), Inches(0.1), Inches(4.5), GOLD)

    # Main title
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "AIエージェント・ストラテジスト"
    run.font.size = Pt(40)
    run.font.color.rgb = GOLD
    run.font.name = FONT_NAME
    run.font.bold = True

    # Subtitle
    tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11.5), Inches(0.8))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = "認定試験 完全対策テキスト"
    run2.font.size = Pt(30)
    run2.font.color.rgb = WHITE
    run2.font.name = FONT_NAME
    run2.font.bold = True

    # Issuer
    tb3 = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(11.5), Inches(0.6))
    tf3 = tb3.text_frame
    p3 = tf3.paragraphs[0]
    run3 = p3.add_run()
    run3.text = "一般社団法人 AICX協会"
    run3.font.size = Pt(22)
    run3.font.color.rgb = LIGHT_GOLD
    run3.font.name = FONT_NAME

    # Details
    details = [
        "初回試験: 2026年7月実施（オンライン・4択ケーススタディ形式）",
        "受験料: ¥14,800（税込）",
        "対象者: DX/AI推進担当者、業務改革担当者、IT企画担当者、コンサルタント、経営層"
    ]

    for i, detail in enumerate(details):
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(4.6 + i * 0.45), Inches(11.5), Inches(0.4))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"• {detail}"
        run.font.size = Pt(16)
        run.font.color.rgb = RGBColor(0xB0, 0xBE, 0xC5)
        run.font.name = FONT_NAME

    return slide

def create_toc_slide(prs):
    """Create table of contents slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide, WHITE)

    # Header bar
    add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.2), NAVY)

    # Gold bottom line of header
    add_rect(slide, Inches(0), Inches(1.2), Inches(13.33), Inches(0.05), GOLD)

    # Title
    add_text_box(slide, "目次", Inches(0.4), Inches(0.2), Inches(12), Inches(0.8),
                font_size=32, font_color=WHITE, bold=True, align=PP_ALIGN.LEFT)

    toc_items = [
        ("A", "業務の構造的理解", "As-Is/To-Be分析、IPO、ECRS、スイムレーン図"),
        ("B", "AIエージェント基礎知識", "生成AI・LLM・プロンプトエンジニアリング"),
        ("C", "RAG（検索拡張生成）", "保存フェーズ・検索生成フェーズ・精度設計"),
        ("D", "AI実装プロセス：実装5Dモデル", "Discovery→Definition→Design→Development→Deployment"),
        ("E", "業務設計・プロセスリエンジニアリング", "業務分解・委任判断・MVP・KPI/ROI"),
        ("F", "AI推進の組織設計", "CoE・人材役割・チェンジマネジメント"),
        ("G", "ガバナンス・倫理・リスク", "ハルシネーション・バイアス・EU AI Act"),
        ("H", "導入・運用・継続改善", "パイロット・ステークホルダー・KPIモニタリング"),
    ]

    col_width = Inches(6.2)
    for i, (letter, title, sub) in enumerate(toc_items):
        col = i % 2
        row = i // 2
        left = Inches(0.4) + col * col_width
        top = Inches(1.5) + row * Inches(1.35)

        # Card background
        add_rect(slide, left, top, col_width - Inches(0.2), Inches(1.2),
                LIGHT_GRAY)

        # Letter badge
        add_rect(slide, left, top, Inches(0.5), Inches(1.2), NAVY)
        add_text_box(slide, letter, left, top + Inches(0.3), Inches(0.5), Inches(0.6),
                    font_size=22, font_color=GOLD, bold=True, align=PP_ALIGN.CENTER)

        # Title
        add_text_box(slide, title, left + Inches(0.6), top + Inches(0.05),
                    col_width - Inches(0.9), Inches(0.5),
                    font_size=16, font_color=NAVY, bold=True)

        # Sub
        add_text_box(slide, sub, left + Inches(0.6), top + Inches(0.55),
                    col_width - Inches(0.9), Inches(0.5),
                    font_size=12, font_color=DARK_GRAY)

    return slide

def create_section_divider(prs, section_letter, section_title, section_subtitle=""):
    """Create a section divider slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide, NAVY)

    # Gold accent bars
    add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.12), GOLD)
    add_rect(slide, Inches(0), Inches(7.38), Inches(13.33), Inches(0.12), GOLD)

    # Large letter
    add_text_box(slide, section_letter, Inches(1), Inches(1.5), Inches(2), Inches(3),
                font_size=120, font_color=RGBColor(0x2A, 0x35, 0xA0), bold=True,
                align=PP_ALIGN.CENTER)

    # Vertical gold bar
    add_rect(slide, Inches(3.2), Inches(1.8), Inches(0.08), Inches(3.5), GOLD)

    # Section title
    add_text_box(slide, f"セクション {section_letter}", Inches(3.5), Inches(2.0), Inches(9), Inches(0.6),
                font_size=18, font_color=GOLD, bold=False)
    add_text_box(slide, section_title, Inches(3.5), Inches(2.6), Inches(9), Inches(1.5),
                font_size=32, font_color=WHITE, bold=True)

    if section_subtitle:
        add_text_box(slide, section_subtitle, Inches(3.5), Inches(4.2), Inches(9), Inches(1.0),
                    font_size=18, font_color=RGBColor(0xB0, 0xBE, 0xC5))

    return slide

def create_content_slide(prs, title, bullets, section_letter=""):
    """Create a content slide with title and bullet points."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide, WHITE)

    # Header bar
    add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.3), NAVY)
    add_rect(slide, Inches(0), Inches(1.3), Inches(13.33), Inches(0.05), GOLD)

    # Section badge
    if section_letter:
        add_rect(slide, Inches(0), Inches(0.15), Inches(0.45), Inches(0.9), GOLD)
        add_text_box(slide, section_letter, Inches(0), Inches(0.2), Inches(0.45), Inches(0.8),
                    font_size=18, font_color=NAVY, bold=True, align=PP_ALIGN.CENTER)

    # Title
    title_left = Inches(0.6) if section_letter else Inches(0.4)
    add_text_box(slide, title, title_left, Inches(0.2), Inches(12.3), Inches(0.9),
                font_size=26, font_color=WHITE, bold=True)

    # Content area
    tf_left = Inches(0.4)
    tf_top = Inches(1.5)
    tf_width = Inches(12.5)

    txBox = slide.shapes.add_textbox(tf_left, tf_top, tf_width, Inches(5.7))
    tf = txBox.text_frame
    tf.word_wrap = True

    first_para = True
    for bullet in bullets:
        if isinstance(bullet, str):
            # Top-level bullet
            if first_para:
                p = tf.paragraphs[0]
                first_para = False
            else:
                p = tf.add_paragraph()

            p.space_before = Pt(6)
            p.space_after = Pt(2)

            run = p.add_run()
            if bullet.startswith("##"):
                # Sub-section header
                run.text = bullet[2:].strip()
                run.font.size = Pt(18)
                run.font.color.rgb = LIGHT_NAVY
                run.font.bold = True
                run.font.name = FONT_NAME
                p.space_before = Pt(10)
            else:
                run.text = f"• {bullet}"
                run.font.size = Pt(18)
                run.font.color.rgb = DARK_GRAY
                run.font.name = FONT_NAME
                run.font.bold = False

        elif isinstance(bullet, tuple):
            # (main_text, [sub_bullets])
            main_text, sub_bullets = bullet

            if first_para:
                p = tf.paragraphs[0]
                first_para = False
            else:
                p = tf.add_paragraph()

            p.space_before = Pt(6)
            run = p.add_run()
            run.text = f"▶ {main_text}"
            run.font.size = Pt(19)
            run.font.color.rgb = NAVY
            run.font.bold = True
            run.font.name = FONT_NAME

            for sub in sub_bullets:
                p2 = tf.add_paragraph()
                p2.space_before = Pt(2)
                p2.level = 1
                run2 = p2.add_run()
                run2.text = f"    ◦ {sub}"
                run2.font.size = Pt(16)
                run2.font.color.rgb = DARK_GRAY
                run2.font.name = FONT_NAME

    return slide

def create_two_column_slide(prs, title, left_header, left_items, right_header, right_items, section_letter=""):
    """Create a two-column content slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide, WHITE)

    # Header bar
    add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.3), NAVY)
    add_rect(slide, Inches(0), Inches(1.3), Inches(13.33), Inches(0.05), GOLD)

    if section_letter:
        add_rect(slide, Inches(0), Inches(0.15), Inches(0.45), Inches(0.9), GOLD)
        add_text_box(slide, section_letter, Inches(0), Inches(0.2), Inches(0.45), Inches(0.8),
                    font_size=18, font_color=NAVY, bold=True, align=PP_ALIGN.CENTER)

    title_left = Inches(0.6) if section_letter else Inches(0.4)
    add_text_box(slide, title, title_left, Inches(0.2), Inches(12.3), Inches(0.9),
                font_size=26, font_color=WHITE, bold=True)

    col_width = Inches(6.0)

    # Left column
    add_rect(slide, Inches(0.3), Inches(1.45), col_width, Inches(0.5), LIGHT_NAVY)
    add_text_box(slide, left_header, Inches(0.4), Inches(1.5), col_width - Inches(0.2), Inches(0.4),
                font_size=16, font_color=WHITE, bold=True)

    left_box = slide.shapes.add_textbox(Inches(0.4), Inches(2.1), col_width - Inches(0.2), Inches(5.0))
    left_tf = left_box.text_frame
    left_tf.word_wrap = True
    first = True
    for item in left_items:
        if first:
            p = left_tf.paragraphs[0]
            first = False
        else:
            p = left_tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        if isinstance(item, tuple):
            run.text = f"▶ {item[0]}"
            run.font.bold = True
            run.font.size = Pt(16)
            run.font.color.rgb = NAVY
            run.font.name = FONT_NAME
            for sub in item[1]:
                p2 = left_tf.add_paragraph()
                r2 = p2.add_run()
                r2.text = f"  ◦ {sub}"
                r2.font.size = Pt(14)
                r2.font.color.rgb = DARK_GRAY
                r2.font.name = FONT_NAME
        else:
            run.text = f"• {item}"
            run.font.size = Pt(16)
            run.font.color.rgb = DARK_GRAY
            run.font.name = FONT_NAME

    # Right column
    right_left = Inches(6.8)
    add_rect(slide, right_left - Inches(0.1), Inches(1.45), col_width, Inches(0.5), GOLD)
    add_text_box(slide, right_header, right_left, Inches(1.5), col_width - Inches(0.2), Inches(0.4),
                font_size=16, font_color=NAVY, bold=True)

    right_box = slide.shapes.add_textbox(right_left, Inches(2.1), col_width - Inches(0.2), Inches(5.0))
    right_tf = right_box.text_frame
    right_tf.word_wrap = True
    first = True
    for item in right_items:
        if first:
            p = right_tf.paragraphs[0]
            first = False
        else:
            p = right_tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        if isinstance(item, tuple):
            run.text = f"▶ {item[0]}"
            run.font.bold = True
            run.font.size = Pt(16)
            run.font.color.rgb = NAVY
            run.font.name = FONT_NAME
            for sub in item[1]:
                p2 = right_tf.add_paragraph()
                r2 = p2.add_run()
                r2.text = f"  ◦ {sub}"
                r2.font.size = Pt(14)
                r2.font.color.rgb = DARK_GRAY
                r2.font.name = FONT_NAME
        else:
            run.text = f"• {item}"
            run.font.size = Pt(16)
            run.font.color.rgb = DARK_GRAY
            run.font.name = FONT_NAME

    return slide

def create_overview_slide(prs):
    """Certification overview slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_background(slide, WHITE)

    add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.3), NAVY)
    add_rect(slide, Inches(0), Inches(1.3), Inches(13.33), Inches(0.05), GOLD)
    add_text_box(slide, "資格概要：AIエージェント・ストラテジスト", Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.9),
                font_size=26, font_color=WHITE, bold=True)

    # Info cards
    cards = [
        ("発行機関", "一般社団法人 AICX協会", NAVY),
        ("初回試験", "2026年7月（オンライン実施）", LIGHT_NAVY),
        ("試験形式", "4択式・ケーススタディ形式", NAVY),
        ("受験料", "¥14,800（税込）", LIGHT_NAVY),
    ]

    card_w = Inches(2.9)
    for i, (label, value, color) in enumerate(cards):
        left = Inches(0.4) + i * Inches(3.1)
        # Card top color bar
        add_rect(slide, left, Inches(1.5), card_w, Inches(0.4), color)
        add_text_box(slide, label, left, Inches(1.52), card_w, Inches(0.36),
                    font_size=14, font_color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        # Card body
        add_rect(slide, left, Inches(1.9), card_w, Inches(0.7), LIGHT_GRAY)
        add_text_box(slide, value, left, Inches(1.95), card_w, Inches(0.65),
                    font_size=15, font_color=NAVY, bold=True, align=PP_ALIGN.CENTER)

    # Target audience
    add_rect(slide, Inches(0.3), Inches(2.85), Inches(12.7), Inches(0.4), GOLD)
    add_text_box(slide, "対象者", Inches(0.5), Inches(2.88), Inches(12.3), Inches(0.35),
                font_size=16, font_color=NAVY, bold=True)

    targets = [
        "DX/AI推進担当者", "業務改革担当者", "IT企画担当者", "コンサルタント", "経営層"
    ]
    target_text = "　　".join([f"✓ {t}" for t in targets])
    add_text_box(slide, target_text, Inches(0.4), Inches(3.35), Inches(12.5), Inches(0.5),
                font_size=16, font_color=DARK_GRAY)

    # 3 Core Domains
    add_rect(slide, Inches(0.3), Inches(4.0), Inches(12.7), Inches(0.45), NAVY)
    add_text_box(slide, "3つのコアスキル領域", Inches(0.5), Inches(4.05), Inches(12.3), Inches(0.4),
                font_size=18, font_color=WHITE, bold=True)

    domains = [
        ("AIエージェント領域", "生成AI・AIエージェントの特性理解、業務適合の自動化設計"),
        ("業務設計領域", "暗黙知の形式知化、ワークフロー再構築"),
        ("組織設計領域", "AIとの協働文化醸成、評価制度の変革"),
    ]

    dom_w = Inches(3.9)
    for i, (dom_title, dom_desc) in enumerate(domains):
        left = Inches(0.4) + i * Inches(4.1)
        add_rect(slide, left, Inches(4.55), dom_w, Inches(0.45), LIGHT_NAVY)
        add_text_box(slide, f"{i+1}. {dom_title}", left + Inches(0.1), Inches(4.58), dom_w - Inches(0.2), Inches(0.4),
                    font_size=14, font_color=WHITE, bold=True)
        add_rect(slide, left, Inches(5.0), dom_w, Inches(1.2), LIGHT_GRAY)
        add_text_box(slide, dom_desc, left + Inches(0.1), Inches(5.05), dom_w - Inches(0.2), Inches(1.1),
                    font_size=13, font_color=DARK_GRAY)

    return slide

# ============================================================
# BUILD THE PRESENTATION
# ============================================================

# Slide 1: Title
create_title_slide(prs)

# Slide 2: Certification Overview
create_overview_slide(prs)

# Slide 3: Table of Contents
create_toc_slide(prs)

# ============================================================
# SECTION A: 業務の構造的理解
# ============================================================
create_section_divider(prs, "A", "業務の構造的理解",
    "As-Is分析・To-Be設計・IPO・ECRS・スイムレーン・暗黙知の形式知化")

create_content_slide(prs, "A-1：As-Is分析とTo-Be設計の基本", [
    ("As-Is分析とは", [
        "現状業務の「ありのまま」を可視化・記録する手法",
        "業務フロー、担当者、所要時間、使用ツール、発生コストを把握",
        "ボトルネック・非効率・重複作業・属人化リスクを特定",
        "インタビュー・ワークショップ・業務観察・ドキュメント分析で実施",
    ]),
    ("To-Be設計とは", [
        "AI導入後の「あるべき姿」を設計する未来の業務フロー",
        "AIが担う業務と人間が担う業務の役割分担を明確化",
        "自動化・効率化により削減できる工数・コストを数値で示す",
        "段階的移行計画（ロードマップ）と合わせて策定",
    ]),
    ("As-Is→To-Beのギャップ分析", [
        "現状と目標の差分（ギャップ）を特定",
        "解決すべき課題の優先順位付けに活用",
        "実装5Dモデル（Discovery/Definition）フェーズで使用",
    ]),
], "A")

create_content_slide(prs, "A-2：IPOフレームワーク（Input-Process-Output）", [
    ("IPOフレームワークとは", [
        "業務を「入力→処理→出力」の3要素で構造化して理解する手法",
        "AIへの業務委任可能性を評価する際の基本ツール",
    ]),
    ("Input（入力）", [
        "業務が始まるために必要なデータ・情報・指示・トリガー",
        "例）顧客からのメール、申請書類、センサーデータ、定期スケジュール",
        "AI委任判断：データの形式・品質・量・入力頻度を評価",
    ]),
    ("Process（処理）", [
        "入力を変換・加工・判断する手順・ロジック・ルール",
        "例）分類判定、情報抽出、計算、コンテンツ生成、承認フロー",
        "AI委任判断：ルール化・パターン化できるかを評価",
    ]),
    ("Output（出力）", [
        "処理の結果として生成されるもの・次のアクション・通知",
        "例）回答メール、レポート、データベース更新、アラート通知",
        "AI委任判断：出力の品質基準・検証方法を設計",
    ]),
    ("AIエージェント設計でのIPO活用", [
        "各業務ステップをIPOで分解→AIが担えるProcessを特定→適切なエージェント設計に繋げる",
    ]),
], "A")

create_content_slide(prs, "A-3：スイムレーン図による業務フロー可視化", [
    ("スイムレーン図とは", [
        "複数の担当者・組織・システムが関与する業務フローを、「泳者のレーン」に見立てて可視化する図",
        "誰が何をいつやるか・どこで情報が受け渡されるかを一目で把握",
    ]),
    ("スイムレーン図の構成要素", [
        "レーン（行/列）：担当者、部門、AIエージェント、外部システムごとに設定",
        "フロー記号：開始・終了（楕円）、処理（矩形）、判断（菱形）、書類（波型矩形）",
        "矢印：業務の流れ・情報の受け渡し方向",
        "クロス線：レーン間の引き継ぎ・連携ポイント",
    ]),
    ("AI導入分析での活用", [
        "現状（As-Is）スイムレーン図：属人化・ボトルネック・手戻りを可視化",
        "将来（To-Be）スイムレーン図：AIエージェントのレーンを追加し役割分担を設計",
        "判断分岐点の洗い出し：AIが自動判断できる条件を明確化",
    ]),
    ("ポイント：AI化に適したプロセス特定", [
        "繰り返し頻度が高い・判断基準が明確・データが構造化されている処理がAI化の候補",
    ]),
], "A")

create_content_slide(prs, "A-4：ECRS原則による業務改善", [
    ("ECRS原則とは", [
        "業務改善・プロセス最適化のための4ステップフレームワーク",
        "AI導入前に「まず無くせないか」を問うことが重要",
    ]),
    ("E：Eliminate（排除）", [
        "その業務を完全に廃止・削除できないかを検討",
        "例）重複レポートの廃止、承認ステップの削除",
        "最も大きな効果。AI化より先に検討すべき",
    ]),
    ("C：Combine（結合）", [
        "複数の業務・ステップを一つにまとめられないかを検討",
        "例）週次報告と月次報告の統合、複数フォームの一本化",
    ]),
    ("R：Rearrange（組み替え）", [
        "業務の順序・担当者・場所・タイミングを変更できないかを検討",
        "例）承認フローの前後入れ替え、担当者変更",
    ]),
    ("S：Simplify（簡素化）", [
        "業務を単純化・標準化・テンプレート化できないかを検討",
        "例）入力フォームの選択式化、チェックリスト化",
        "AIはSimplify後の業務を自動化するとより高い効果",
    ]),
    ("試験ポイント：ECRS後にAI化", [
        "まずECRSで業務最適化→残った標準業務をAIで自動化が正しい順序",
    ]),
], "A")

create_content_slide(prs, "A-5：暗黙知の形式知化", [
    ("暗黙知と形式知の定義", [
        "暗黙知（Tacit Knowledge）：言語化・文書化されていない経験・勘・コツ・判断基準",
        "形式知（Explicit Knowledge）：文書・マニュアル・データベースとして共有可能な知識",
        "野中郁次郎のSECIモデル：知識創造の4プロセス（共同化・表出化・連結化・内面化）",
    ]),
    ("AIへの業務委任に暗黙知形式知化が必要な理由", [
        "AIは形式知のみ学習・活用できる→暗黙知はAIに教えられない",
        "プロンプト・RAGナレッジベース・判断ルールの材料となる",
        "組織の属人化リスク解消・ナレッジ継承にも貢献",
    ]),
    ("形式知化の手法", [
        "インタビュー法：熟練者へのヒアリング・思考発話法",
        "業務観察法：実際の作業を観察・記録・分析",
        "ドキュメント分析：既存マニュアル・議事録・メールから抽出",
        "ケーススタディ整理：過去の判断事例をIF-THEN形式で整理",
        "フローチャート化：判断分岐を図示→AIのロジック設計に活用",
    ]),
    ("RAGとの連携", [
        "形式知化した知識をドキュメント化→RAGのナレッジベースへ登録",
        "AIエージェントが参照できる「業務知識DB」を構築",
    ]),
], "A")

# ============================================================
# SECTION B: AIエージェント基礎知識
# ============================================================
create_section_divider(prs, "B", "AIエージェント基礎知識",
    "生成AI・LLM・プロンプトエンジニアリング・エージェント類型")

create_content_slide(prs, "B-1：生成AIとLLMの基礎", [
    ("大規模言語モデル（LLM）とは", [
        "大量のテキストデータで学習した深層学習モデル",
        "テキスト生成・翻訳・要約・質問応答・コード生成が可能",
        "代表例：GPT-4/4o、Claude、Gemini、Llama等",
    ]),
    ("トークン（Token）", [
        "LLMがテキストを処理する最小単位（単語・部分文字列・文字）",
        "日本語は英語より1文字あたり多くのトークンを消費する傾向",
        "コスト・速度・コンテキスト容量はトークン数で決まる",
    ]),
    ("コンテキストウィンドウ", [
        "LLMが一度に処理できるトークンの最大量",
        "大きいほど長い文書・会話履歴を参照できる",
        "例）GPT-4：128K tokens、Claude 3：200K tokens",
        "超過した場合は古い情報が切り捨てられる（重要な設計考慮点）",
    ]),
    ("温度パラメータ（Temperature）", [
        "出力のランダム性・創造性を制御するパラメータ（0.0〜2.0）",
        "低い値（0.0〜0.3）：決定論的・一貫性重視→業務処理・正確さが求められるタスクに適切",
        "高い値（0.7〜1.5）：創造的・多様な出力→アイデア出し・文章生成に活用",
        "ビジネス用途では一般的に低い温度設定を推奨",
    ]),
], "B")

create_two_column_slide(prs, "B-2：生成AI vs 従来型AI の比較",
    "従来型AI（ルールベース・ML）", [
        ("特徴", ["特定タスク向けに設計", "大量の教師データが必要", "明示的なルール・パターン学習"]),
        ("強み", ["精度が高い（特化タスク）", "動作の予測可能性が高い", "処理速度が速い"]),
        ("弱み", ["汎用性が低い", "新しいタスクには再学習が必要", "柔軟な対話・生成が困難"]),
        ("例", ["画像認識、スパムフィルタ、推薦エンジン", "異常検知、需要予測"]),
    ],
    "生成AI（LLMベース）", [
        ("特徴", ["汎用的な言語理解・生成能力", "プロンプトで動作変更可能", "ゼロショット・少数ショット学習"]),
        ("強み", ["高い汎用性・柔軟性", "自然言語での指示が可能", "マルチタスク対応"]),
        ("弱み", ["ハルシネーション（幻覚）リスク", "最新情報の欠如（カットオフ）", "高コスト・遅延"]),
        ("例", ["文書生成、対話、要約、翻訳", "コード生成、分析支援"]),
    ],
    "B"
)

create_content_slide(prs, "B-3：AIエージェントの定義と特性", [
    ("AIエージェントとは", [
        "目標を与えられると、自律的に計画を立て、ツール・APIを使い、行動を繰り返しながら目標達成を目指すAIシステム",
        "単純な「質問→回答」の生成AIとは異なり、「計画→実行→観察→修正」のループで動作",
    ]),
    ("自律性（Autonomy）", [
        "人間の逐一指示なしに、自らの判断でアクションを選択・実行",
        "外部ツール（検索、コード実行、API）の呼び出しを自ら判断",
    ]),
    ("目標指向性（Goal-Oriented）", [
        "与えられた最終目標に向けてサブゴールを分解・順次達成",
        "ReAct（Reasoning + Acting）パターン：推論と行動を交互に実行",
    ]),
    ("ループ処理（Agentic Loop）", [
        "Observe（観察）→ Think（推論・計画）→ Act（実行）→ Observe...のサイクル",
        "エラーや予期しない結果に対して自動的に軌道修正",
        "Multi-Agent：複数エージェントが協調してタスクを分担",
    ]),
    ("AIエージェントの主要コンポーネント", [
        "LLM（脳）：推論・計画・言語理解",
        "ツール（手）：Web検索・コード実行・DB操作・API連携",
        "メモリ（記憶）：短期（会話履歴）・長期（ベクトルDB）",
        "オーケストレーター（指揮者）：エージェント間の調整",
    ]),
], "B")

create_content_slide(prs, "B-4：AIエージェントの類型", [
    ("①カスタム型チャットボット（GPTs・Gems等）", [
        "特定の指示・ナレッジ・ペルソナを持つカスタムAIアシスタント",
        "OpenAI GPTs、Google Gems、Claude Projects等で構築",
        "システムプロンプト・知識ファイル・ツール権限を設定",
        "適用場面：社内FAQ、特定ドメイン質問応答、コンシェルジュ",
        "限界：単一会話内での対応・複雑な外部連携には制限",
    ]),
    ("②ワークフロー型AIエージェント", [
        "事前定義されたフローに沿って複数のAI処理・ツール呼び出しを自動実行",
        "n8n、Zapier、Make、Azure Logic Appsなどで構築",
        "トリガー（起動条件）→ステップ（処理）→条件分岐→アクション",
        "適用場面：定型業務自動化、データパイプライン、通知フロー",
        "特徴：予測可能・安定・監視しやすい（ガバナンス管理が容易）",
    ]),
    ("③汎用型AIエージェント", [
        "高い自律性を持ち、複雑な非定型タスクを動的に計画・実行",
        "LangChain、LangGraph、AutoGen、CrewAI等で構築",
        "ツール選択・サブタスク分解・エラー回復を自律的に実施",
        "適用場面：調査・分析・複数システム横断処理・意思決定支援",
        "特徴：強力だが監視・ガバナンス設計が重要",
    ]),
], "B")

create_content_slide(prs, "B-5：プロンプトエンジニアリング", [
    ("プロンプトエンジニアリングとは", [
        "LLMから望ましい出力を引き出すための入力文（プロンプト）設計技術",
        "AI活用の質は「どう指示するか」で大きく変わる",
    ]),
    ("基本構成要素", [
        "役割（Role）：AIに演じさせるペルソナ「あなたは〇〇の専門家です」",
        "文脈（Context）：背景情報・目的・制約条件の提供",
        "指示（Instruction）：具体的な実行タスク・形式・出力長の指定",
        "例示（Examples）：期待する入出力の例を提示",
    ]),
    ("Few-shot プロンプティング", [
        "入出力例を複数示してLLMにパターンを学習させる手法",
        "例）「以下の形式でメールを分類してください→例1...例2...→本番データ」",
        "ゼロショット（例なし）より精度向上。コスト増とのトレードオフ",
    ]),
    ("Chain-of-Thought（CoT）プロンプティング", [
        "「ステップごとに考えてください」と推論プロセスを明示させる手法",
        "複雑な論理問題・計算・多段階判断で精度向上",
        "例）「まず〇〇を確認し、次に△△を検討し、最後に結論を出してください」",
    ]),
    ("業務プロンプト設計のベストプラクティス", [
        "具体的・明確・構造化された指示を心がける",
        "出力フォーマット（JSON・箇条書き・テーブル等）を明示",
        "制約条件（文字数・禁止事項・トーン）を含める",
        "バージョン管理・A/Bテストで継続的に改善",
    ]),
], "B")

# ============================================================
# SECTION C: RAG
# ============================================================
create_section_divider(prs, "C", "RAG（検索拡張生成）",
    "保存フェーズ・検索生成フェーズ・RAG精度設計")

create_content_slide(prs, "C-1：RAGの定義と目的", [
    ("RAG（Retrieval-Augmented Generation）とは", [
        "LLMの回答生成時に、外部知識ベースから関連情報を動的に検索し、コンテキストとして注入する技術",
        "「検索（Retrieval）」と「生成（Generation）」を組み合わせたアーキテクチャ",
    ]),
    ("RAGが解決する3つの課題", [
        "①ハルシネーション（幻覚）の抑制：根拠となる文書を参照させることで事実と異なる生成を抑制",
        "②最新性の確保：LLMの学習データカットオフ後の情報もリアルタイム参照可能",
        "③根拠・引用の提示：「どの文書に基づいているか」を明示し信頼性・透明性を確保",
    ]),
    ("RAGが特に有効なユースケース", [
        "社内規程・製品マニュアル・法令に基づく質問応答",
        "最新ニュース・市場情報・業績データを使った分析支援",
        "カスタマーサポート・ヘルプデスクの自動応答",
    ]),
    ("RAGとファインチューニングの使い分け", [
        "RAG：知識の追加・更新が頻繁、引用が重要、コスト重視の場合",
        "ファインチューニング：特定スタイル・形式への適応、推論時コスト削減が目的",
    ]),
], "C")

create_content_slide(prs, "C-2：RAG 保存フェーズ（インデックス構築）", [
    ("STEP 1：ドキュメント分割（チャンキング）", [
        "大きな文書を検索しやすい小さなチャンク（断片）に分割",
        "固定サイズ分割：一定文字数/トークン数で機械的に分割（シンプルだが文脈が途切れるリスク）",
        "意味単位分割：段落・セクション・文書構造に基づいて分割（精度向上）",
        "チャンクサイズの設計が精度に大きく影響（詳細はC-4）",
    ]),
    ("STEP 2：埋め込みベクトル化（Embedding）", [
        "テキストを多次元ベクトル（数値配列）に変換する処理",
        "埋め込みモデル：OpenAI text-embedding-3、Cohere Embed等",
        "意味的に近いテキストはベクトル空間上で近い位置に配置される",
        "日本語対応モデルの選定が精度のカギ",
    ]),
    ("STEP 3：ベクトルデータベースへの格納", [
        "変換したベクトルを高速検索に最適化したDBに保存",
        "代表的なベクトルDB：Pinecone、Weaviate、Chroma、pgvector",
        "元テキスト・メタデータ（ファイル名・作成日・カテゴリ）もセットで保存",
        "メタデータはフィルタリング・ランキング精度向上に活用（詳細はC-4）",
    ]),
], "C")

create_content_slide(prs, "C-3：RAG 検索・生成フェーズ（クエリ処理）", [
    ("STEP 4：クエリ処理（Query Processing）", [
        "ユーザーの質問（クエリ）をEmbeddingモデルでベクトル化",
        "クエリ拡張：同義語追加・分解・仮説文書生成（HyDE）で検索精度向上",
        "多言語クエリ：日本語←→英語の変換・統一化",
    ]),
    ("STEP 5：類似度検索（Semantic Search）", [
        "クエリベクトルとチャンクベクトルの類似度を計算",
        "コサイン類似度・内積・ユークリッド距離で近いチャンクを特定",
        "Top-K件のチャンクを候補として取得",
        "ハイブリッド検索：ベクトル検索＋キーワード検索（BM25）の組み合わせで精度向上",
    ]),
    ("STEP 6：コンテキスト注入と生成（Context Injection & Generation）", [
        "取得したチャンクをLLMのプロンプトに「コンテキスト」として挿入",
        "LLMはコンテキストを参照しながら回答を生成",
        "引用情報（チャンク出典・ページ数）を合わせて提示",
        "「与えられた文書のみに基づいて回答してください」という制約プロンプトが重要",
    ]),
    ("リランキング（Re-ranking）", [
        "初期検索で取得したK件を、より精密なモデルで再順位付け",
        "Cohere Rerank、Cross-Encoder等を使用",
        "精度は向上するがレイテンシ・コストが増加",
    ]),
], "C")

create_content_slide(prs, "C-4：RAG精度設計のポイント", [
    ("チャンクサイズの最適化", [
        "小さなチャンク（128〜256トークン）：精密検索・高精度だが文脈が不足しがち",
        "大きなチャンク（512〜1024トークン）：豊富な文脈だがノイズ増・精度低下リスク",
        "Parent-Child chunking：小チャンクで検索→親チャンクで文脈を返す高度手法",
        "オーバーラップ：チャンク間で一部テキストを重複させ文脈の断絶を防ぐ",
    ]),
    ("リランキング（Re-ranking）", [
        "ベクトル検索は意味的近さで検索するが関連性が完全ではない",
        "Cross-Encoderモデルでクエリとチャンクのペアを評価・再順位付け",
        "精度向上と計算コストのトレードオフを考慮して設計",
    ]),
    ("メタデータ活用", [
        "チャンクに文書名・作成日・部門・バージョン等のメタデータを付与",
        "フィルタリング：「2024年以降の文書のみ」「営業部門のマニュアルのみ」",
        "メタデータによる前フィルタ→ベクトル検索の精度・速度向上",
    ]),
    ("RAG評価指標", [
        "Faithfulness：回答が検索文書に忠実か（幻覚の有無）",
        "Answer Relevancy：回答がクエリに対して適切か",
        "Context Precision/Recall：適切なチャンクが検索されているか",
        "RAGASフレームワーク：上記指標の自動評価ツール",
    ]),
], "C")

# ============================================================
# SECTION D: 実装5Dモデル
# ============================================================
create_section_divider(prs, "D", "AI実装プロセス：実装5Dモデル",
    "Discovery→Definition→Design→Development→Deployment")

create_content_slide(prs, "D-1：実装5Dモデル 概要", [
    ("実装5Dモデルとは", [
        "AI・AIエージェント導入を成功させるための5段階の実装フレームワーク",
        "各フェーズに明確な目的・成果物・参加者・判断基準を設定",
        "ウォーターフォール型ではなく、各フェーズで検証と学習を繰り返すアジャイル的アプローチ",
    ]),
    ("5Dモデルのフェーズ一覧", [
        "D1：Discovery（発見）　　→ 課題発見・ユースケース特定",
        "D2：Definition（定義）　 → 要件定義・成功指標設定",
        "D3：Design（設計）　　　 → 業務フロー・エージェント設計",
        "D4：Development & PoC　 → 開発・概念実証",
        "D5：Deployment & Scale　→ 展開・スケーリング",
    ]),
    ("各フェーズの重要性", [
        "D1-D2をスキップすると「技術ありきの導入」になり失敗リスク大",
        "D4のPoCで早期に実現可能性・ROIを検証することが成功の鍵",
        "D5では変化管理・継続的改善の仕組みを整備",
    ]),
], "D")

create_content_slide(prs, "D-2：Discovery（課題発見・ユースケース特定）", [
    ("Discovery フェーズの目的", [
        "解決すべきビジネス課題を特定し、AIで対処すべきユースケースを発見する",
        "「AIが使いたい」ではなく「課題解決にAIが最適か」を問う",
    ]),
    ("主な活動", [
        "ステークホルダーインタビュー：経営層・業務担当者・現場からの課題ヒアリング",
        "As-Is業務分析：現状業務フロー・ボトルネック・工数・コストの可視化",
        "ペインポイントマッピング：頻度×深刻度で課題を優先順位付け",
        "ユースケース候補リスト化：AIで解決できる可能性の高い課題を列挙",
    ]),
    ("ユースケース評価基準（AIに適した業務の特徴）", [
        "繰り返し頻度が高い（週次・日次以上）",
        "判断基準がルール化・データ化できる",
        "大量データの処理を伴う",
        "人間が行う際に時間・品質のばらつきが大きい",
        "失敗コスト（可逆性）が許容範囲内",
    ]),
    ("成果物", [
        "課題マップ（ペインポイント一覧）・優先ユースケースリスト・初期ROI試算",
    ]),
], "D")

create_content_slide(prs, "D-3：Definition（要件定義・成功指標設定）", [
    ("Definition フェーズの目的", [
        "選定したユースケースの要件を詳細定義し、成功/失敗の判断基準を設定する",
    ]),
    ("要件定義の内容", [
        "機能要件：AIが実行すべき具体的な機能・入出力・処理ロジック",
        "非機能要件：応答時間・精度・可用性・セキュリティ・スケーラビリティ",
        "データ要件：必要なデータの種類・量・品質・取得元・更新頻度",
        "インテグレーション要件：連携が必要な既存システム・API・DB",
    ]),
    ("成功指標（KPI）の設定", [
        "定量指標：処理時間削減率（例：80%削減）、コスト削減額、精度/F1スコア",
        "定性指標：ユーザー満足度、従業員の業務負荷軽減感",
        "ベースライン設定：現状値を計測し変化を測定できる状態にする",
        "測定方法・報告サイクルを事前に決定",
    ]),
    ("ROI（投資対効果）の設計", [
        "投資コスト：開発費・ライセンス料・インフラ費・運用費・教育費",
        "効果：人件費削減・収益増加・品質向上・リスク低減",
        "ROI = (効果 - コスト) / コスト × 100%",
        "回収期間（ペイバックピリオド）も合わせて試算",
    ]),
], "D")

create_content_slide(prs, "D-4：Design（業務フロー・エージェント設計）", [
    ("Design フェーズの目的", [
        "To-Be業務フローとAIエージェントの詳細設計を行う",
    ]),
    ("To-Be業務フロー設計", [
        "スイムレーン図でAIと人間の役割分担を可視化",
        "AIが処理するステップと人間がレビュー・承認するチェックポイントを明確化",
        "例外処理フロー：AIが対応できないケースの人間へのエスカレーション設計",
    ]),
    ("AIエージェント設計の要素", [
        "エージェントタイプ選択：チャットボット型/ワークフロー型/汎用型",
        "ツール設計：使用するAPI・DB・外部サービスの定義",
        "プロンプト設計：システムプロンプト・Few-shot例・制約条件",
        "メモリ設計：会話履歴の保持範囲・RAGナレッジベース",
        "エラーハンドリング：失敗時の動作・リトライ・通知設計",
    ]),
    ("ガードレール設計", [
        "出力フィルタリング：不適切・危険な出力を防止するセーフガード",
        "Human-in-the-Loop：重要な判断への人間の確認ステップ",
        "監査ログ：AIの判断・行動履歴の記録",
        "ロールバック計画：問題発生時の切り戻し手順",
    ]),
], "D")

create_content_slide(prs, "D-5：Development & PoC（開発・概念実証）", [
    ("Development & PoC フェーズの目的", [
        "最小限の機能でプロトタイプを開発し、実現可能性・ROIを早期検証する",
    ]),
    ("PoC（概念実証）の設計原則", [
        "MVP思想：必要最小限の機能のみで検証（機能を絞ることが重要）",
        "実際の業務データ・ユーザーで検証（サンドボックス環境で）",
        "検証期間：2〜8週間が目安（長引くと「PoC地獄」に陥るリスク）",
        "判断基準を事前に設定（GoかNo-Goの明確な閾値）",
    ]),
    ("MVP（最小実行可能プロダクト）設計", [
        "コアユースケースの一つに絞る",
        "既存ツール・API・ノーコード/ローコードを最大活用",
        "精度・速度よりも「ユースケースの実現可能性」を最優先で検証",
        "エンドユーザーへの早期提示でフィードバック収集",
    ]),
    ("PoCから本番への判断基準", [
        "精度指標が事前設定閾値を超えているか",
        "ユーザーが実際に使えるか（UX・操作性）",
        "ROI見込みが期待値を満たしているか",
        "セキュリティ・コンプライアンス要件を満たしているか",
    ]),
], "D")

create_content_slide(prs, "D-6：Deployment & Scale（展開・スケーリング）", [
    ("Deployment & Scale フェーズの目的", [
        "PoCを本番環境に展開し、継続的に改善・スケールさせる",
    ]),
    ("展開戦略", [
        "パイロット展開：特定部門・チーム・ユースケースから開始（リスク限定）",
        "段階的ロールアウト：成功を確認しながら展開範囲を徐々に拡大",
        "フルデプロイ：全社・全ユーザーへの展開",
    ]),
    ("スケーリング設計", [
        "インフラスケール：処理量増加に対応したAPIレート制限・コスト管理",
        "ユースケーススケール：1つの成功事例から横展開・類似ユースケースへの適用",
        "組織スケール：使用部門・地域・拠点の拡大",
    ]),
    ("継続的改善サイクル", [
        "KPIモニタリング：定期的な精度・効率・満足度測定",
        "フィードバックループ：ユーザーフィードバック→プロンプト改善→モデル更新",
        "定期レビュー：月次・四半期でのKPI評価・改善計画策定",
        "モデルアップデート対応：LLMのバージョンアップへの追従・再評価",
    ]),
], "D")

# ============================================================
# SECTION E: 業務設計・プロセスリエンジニアリング
# ============================================================
create_section_divider(prs, "E", "業務設計・プロセスリエンジニアリング",
    "業務分解・AI委任判断・MVP・KPI/ROI設計")

create_content_slide(prs, "E-1：業務分解の手法", [
    ("業務分解とは", [
        "大きな業務・プロセスをAI化・分析できる単位に細分化する手法",
        "WBS（Work Breakdown Structure）的な思想で業務を階層的に分解",
    ]),
    ("業務分解の3レベル", [
        "レベル1（業務領域）：例）「顧客対応業務」「受発注業務」",
        "レベル2（業務プロセス）：例）「問い合わせ受付」「在庫確認」「発注処理」",
        "レベル3（業務タスク）：例）「メール本文の読み取り」「在庫DBの照会」「数量の入力」",
    ]),
    ("タスクレベルでのAI化評価", [
        "入力データの形式：テキスト・数値・画像・音声",
        "処理の複雑さ：ルール明確/曖昧、判断の裁量幅",
        "出力の検証しやすさ：正誤が明確か、人間によるレビューが必要か",
        "発生頻度と件数：自動化ROIを左右する重要指標",
    ]),
    ("業務分解の実施手順", [
        "①現状業務のヒアリング・ドキュメント収集",
        "②主要プロセスの洗い出しと順序整理",
        "③各プロセスのタスクレベルへの細分化",
        "④IPOフレームワークで各タスクを構造化",
        "⑤AIデリゲーション評価（次スライド参照）",
    ]),
], "E")

create_content_slide(prs, "E-2：AIへの業務委任判断基準（デリゲーション）", [
    ("AIデリゲーション（AI委任）とは", [
        "人間が行っている業務タスクをAIに委任する判断のフレームワーク",
        "「何をAIに任せ、何を人間が担うか」を体系的に決定する",
    ]),
    ("委任判断の4象限マトリクス", [
        "高頻度×高ルール化　→ 完全AI委任（自動化優先）",
        "高頻度×低ルール化　→ AI支援（人間が最終判断・AIが草案）",
        "低頻度×高ルール化　→ 条件付きAI委任（閾値設定・人間確認）",
        "低頻度×低ルール化　→ 人間主体（AIは参考意見のみ）",
    ]),
    ("委任適性評価の6基準", [
        "①繰り返し性：同じパターンが繰り返されるか",
        "②データ依存性：判断がデータ・テキストに基づくか",
        "③ルール明確性：判断基準を明文化できるか",
        "④影響度・可逆性：誤判断した場合の影響範囲・修正の容易さ",
        "⑤説明責任：AIの判断で法的・倫理的問題が生じないか",
        "⑥ユーザー受容性：関係者・顧客がAI対応を受け入れるか",
    ]),
    ("Human-in-the-Loop設計", [
        "全面委任ではなくAI判断に人間の確認を挟む設計",
        "確認が必要なトリガー条件（信頼度スコア閾値等）を設定",
    ]),
], "E")

create_content_slide(prs, "E-3：As-Is業務分析とTo-Be業務設計", [
    ("As-Is業務分析のポイント（ペインポイント特定）", [
        "定量的分析：各タスクの所要時間・件数・担当人数・エラー率を計測",
        "定性的分析：担当者が感じる不満・属人化・判断の迷いを抽出",
        "ペインポイントの分類：効率性の問題・品質の問題・属人性の問題・コストの問題",
        "ロードマップ優先度付け：影響度×AI化可能性で優先順位を決定",
    ]),
    ("よくあるペインポイント事例", [
        "手動コピペ・転記作業（時間浪費・ミス多発）",
        "同じ質問への繰り返し回答（FAQ対応）",
        "大量文書のスキャン・内容確認・要点抽出",
        "定型レポートの作成・集計・グラフ化",
        "複数システムへの同一データの二重入力",
    ]),
    ("To-Be業務設計の原則", [
        "AIファースト設計：「まずAIが処理し、必要に応じて人間が介入」の発想",
        "シームレスな引き継ぎ：AI→人間へのエスカレーションを滑らかに設計",
        "例外処理の明確化：AIが対応できないケースの処理フロー",
        "継続改善サイクルの組み込み：フィードバックを業務フローに組み入れる",
    ]),
], "E")

create_content_slide(prs, "E-4：MVP設計とKPI・ROI設計", [
    ("MVP（最小実行可能プロダクト）設計", [
        "最も価値が高く実証しやすいユースケース一つに機能を絞る",
        "ゴールは「完璧な製品」ではなく「仮説検証」",
        "2〜4週間でデモできるレベルを目指す",
        "フィードバックを元にイテレーションを繰り返す",
    ]),
    ("KPI設計の原則", [
        "SMARTな目標：Specific（具体的）・Measurable（計測可能）・Achievable（達成可能）・Relevant（関連性）・Time-bound（期限）",
        "結果KPI：業務効率化率・コスト削減額・エラー率・顧客満足度（NPS）",
        "先行KPI：AI利用率・プロセス完了時間・ユーザーフィードバック",
        "ベースラインの事前計測が必須",
    ]),
    ("ROI（投資対効果）設計", [
        "コスト項目：初期開発費・APIライセンス・インフラ費・保守運用費・教育費",
        "効果項目（定量）：人件費削減（工数×単価）・売上増加・ミス削減コスト",
        "効果項目（定性）：ブランド価値・従業員満足度・ナレッジ蓄積",
        "計算式：ROI = (純効果 ÷ 総投資コスト) × 100%",
        "回収期間（Break-even Point）も経営層向けに試算して提示",
    ]),
], "E")

# ============================================================
# SECTION F: AI推進の組織設計
# ============================================================
create_section_divider(prs, "F", "AI推進の組織設計",
    "CoE・人材役割・組織モデル・チェンジマネジメント")

create_content_slide(prs, "F-1：AI推進の組織モデル", [
    ("①中央集権型（CoE主導モデル）", [
        "AI推進専門部署（CoE）が全社のAI戦略・実装・ガバナンスを一元管理",
        "メリット：標準化・ノウハウ集約・一貫したガバナンス・規模の経済",
        "デメリット：ビジネス部門との距離感・リソース競合・ボトルネック化リスク",
        "向いている組織：規制産業（金融・医療）・AI初期段階・標準化を重視する大企業",
    ]),
    ("②分散型（各部門主体モデル）", [
        "各事業部門が独自にAIプロジェクトを推進・実装",
        "メリット：ビジネス課題への密着・迅速な意思決定・現場主導",
        "デメリット：ツール乱立・品質不均一・ノウハウ断絶・ガバナンス困難",
        "向いている組織：変化の速い業界・部門間独立性が高い・AI成熟度が高い企業",
    ]),
    ("③ハイブリッド型（連邦型モデル）", [
        "CoEが標準・ガバナンス・共通基盤を提供し、各部門が実装を主体的に推進",
        "メリット：標準化と現場密着を両立・スケーラブル・ベストプラクティス共有",
        "デメリット：役割分担の曖昧さ・調整コスト",
        "向いている組織：大企業・グローバル展開・AI成熟フェーズの企業（最も推奨）",
    ]),
], "F")

create_content_slide(prs, "F-2：CoE（Center of Excellence）の役割と機能", [
    ("CoEとは", [
        "AI・DXを推進するための専門知識・ベストプラクティス・標準を集約した組織横断的な専門チーム",
        "「推進母体」としての機能と「支援・共有」機能を両立",
    ]),
    ("CoEの主要機能", [
        "戦略立案：全社AIロードマップ・優先ユースケースの決定",
        "標準化：技術スタック・開発ガイドライン・プロンプト標準・セキュリティポリシー",
        "能力開発：社内AI教育・研修・資格取得支援・コミュニティ運営",
        "ガバナンス：倫理審査・リスク評価・コンプライアンス確認",
        "PoC支援：各部門のAIプロジェクトへの技術支援・ベストプラクティス共有",
        "ベンダー管理：AIベンダー・ツールの評価・調達・契約管理",
    ]),
    ("CoEの理想的な構成メンバー", [
        "AIアーキテクト/エンジニア：技術設計・実装支援",
        "データサイエンティスト：モデル評価・データ分析",
        "AI倫理/ガバナンス担当：リスク・コンプライアンス",
        "業務改革担当（BPR）：業務プロセス設計支援",
        "プロジェクトマネージャー：推進管理・調整",
        "チェンジマネジメント担当：組織変革・教育・コミュニケーション",
    ]),
], "F")

create_content_slide(prs, "F-3：AI推進の人材役割", [
    ("スポンサー（Sponsor）", [
        "役割：AI推進プロジェクトに権限・予算・政治的支援を与える経営幹部",
        "重要性：組織変革を成功させる最重要要素。スポンサーシップ不在は失敗の最大要因",
        "行動：定期的な関与・現場への優先度メッセージ発信・リソース確保",
    ]),
    ("推進者（Champion/Change Agent）", [
        "役割：現場とCoEを繋ぐ部門内のAI推進リーダー",
        "特性：業務知識と技術への興味・コミュニケーション能力・変化への意欲",
        "行動：部門内の課題発見・PoC推進・現場への普及・フィードバック収集",
        "別名：AI Champion、DX推進担当、デジタルアンバサダー",
    ]),
    ("実装者（Implementer）", [
        "役割：AIシステムを実際に構築・設定・保守するエンジニア・開発者",
        "スキル：LLMアプリ開発・プロンプトエンジニアリング・API連携・RAG構築",
        "社内育成 vs 外部調達：初期は外部活用し並行して内製化を進める戦略が一般的",
    ]),
    ("エンドユーザー（End User）", [
        "役割：AIツールを日常業務で使用する最終利用者",
        "重要性：利用率と活用品質がROIを左右する",
        "支援策：適切なトレーニング・サポート体制・フィードバックチャンネルの整備",
    ]),
], "F")

create_content_slide(prs, "F-4：チェンジマネジメント", [
    ("チェンジマネジメントとは", [
        "AI導入に伴う組織・プロセス・文化の変革を計画的に管理し、人間側の適応を促す活動",
        "技術的実装と同等かそれ以上に重要。多くのAI導入失敗は技術ではなく人・組織の問題",
    ]),
    ("コッターの変革8段階モデルのAI導入への適用", [
        "①危機感の醸成：「なぜ今AIか」を明確に伝える",
        "②推進連合の形成：スポンサー・推進者チームの構築",
        "③ビジョンと戦略の策定：AI活用のあるべき姿と実行計画",
        "④ビジョンの周知：全社コミュニケーション・説明会・Q&A",
        "⑤行動のエンパワメント：障害除去・トレーニング提供",
        "⑥短期成果の創出：小さな成功（クイックウィン）を可視化",
        "⑦改善の定着：成功事例の横展開・定例化",
        "⑧文化への定着：AI活用を評価制度・採用基準に組み込む",
    ]),
    ("AI導入における典型的な抵抗要因と対策", [
        "「仕事を奪われる」恐怖 → AIは代替でなく「拡張」であることを丁寧に伝える",
        "「使いこなせない」不安 → 実践的トレーニング・サポート体制",
        "「今まで通りで十分」慣性 → 現状維持のコスト・リスクを可視化",
    ]),
], "F")

# ============================================================
# SECTION G: ガバナンス・倫理・リスク
# ============================================================
create_section_divider(prs, "G", "ガバナンス・倫理・リスク",
    "ハルシネーション・バイアス・プライバシー・XAI・EU AI Act")

create_content_slide(prs, "G-1：AIガバナンスの基本原則", [
    ("AIガバナンスとは", [
        "AI技術の開発・導入・運用を適切に管理・監督するための方針・プロセス・体制の総称",
        "技術的安全性と倫理的整合性を担保しながらビジネス価値を最大化する",
    ]),
    ("AIガバナンスの5大原則", [
        "①人間中心性：AIは人間の監督下に置かれ、最終判断は人間が行う",
        "②透明性・説明可能性：AIがなぜその判断をしたか説明できる",
        "③公平性・非差別：性別・人種・年齢等によるバイアスを排除する",
        "④安全性・セキュリティ：悪用・誤作動・情報漏洩から保護する",
        "⑤説明責任：AIシステムの運用者・開発者が結果に責任を持つ",
    ]),
    ("企業AIガバナンスの実践要素", [
        "AIポリシー（利用規程）：使用可能なAIツール・禁止事項・データ取扱い基準",
        "リスク評価プロセス：新規AI導入前の倫理・法的リスク審査",
        "監査・モニタリング：AIの動作・判断を継続的に記録・検証",
        "インシデント対応計画：問題発生時の対処・エスカレーション手順",
    ]),
], "G")

create_content_slide(prs, "G-2：ハルシネーション（幻覚）リスク", [
    ("ハルシネーションとは", [
        "LLMが事実と異なる情報を自信を持って生成してしまう現象",
        "モデルが学習データから確率的にテキストを生成するために発生",
        "ビジネス利用における最大のリスクの一つ",
    ]),
    ("ハルシネーションの種類", [
        "事実的誤り：存在しない法律・数値・人物・出来事を生成",
        "文脈的誤り：正確な情報だが質問の文脈と合わない回答",
        "内部矛盾：同一回答内で矛盾する記述が存在",
    ]),
    ("ハルシネーション低減策", [
        "RAG（検索拡張生成）：根拠文書を参照させ事実に基づいた回答を生成",
        "低温度パラメータ：ランダム性を下げ一貫性を確保",
        "ファクトチェック指示：「不確実な場合は『わかりません』と回答して」",
        "グラウンディング：「以下の文書のみに基づいて回答してください」",
        "出力検証：別のAI/人間によるファクトチェックレイヤー",
    ]),
    ("業務別リスク評価", [
        "高リスク：医療診断・法的判断・財務報告 → 必ず人間の最終確認",
        "中リスク：顧客対応・マーケティング文書 → レビュープロセス設計",
        "低リスク：ドラフト作成・内部メモ → AIの判断を活用しやすい",
    ]),
], "G")

create_content_slide(prs, "G-3：バイアスとフェアネス", [
    ("AIバイアスとは", [
        "AIシステムが特定グループに対して不公平な判断・出力をする傾向",
        "学習データに含まれる社会的偏見・歴史的不均衡がモデルに引き継がれる",
    ]),
    ("バイアスの主要類型", [
        "データバイアス：学習データに特定グループのデータが過剰/過少に含まれる",
        "アルゴリズムバイアス：モデル設計・目的関数がバイアスを増幅させる",
        "測定バイアス：特定グループで精度が低い特徴量の使用",
        "フィードバックバイアス：バイアスのある結果が再学習データとなり悪化",
    ]),
    ("バイアス検出と対策", [
        "代表性チェック：学習データの人口統計的分布の確認",
        "公平性指標評価：グループ間の精度・偽陽性率・偽陰性率の差を計測",
        "カウンターファクチュアル検証：属性のみ変えたテストケースで偏りを確認",
        "Fairness-aware学習：公平性制約を組み込んだモデル設計",
    ]),
    ("採用・審査業務でのバイアスリスク", [
        "採用AIが過去データの性別・学歴バイアスを学習するリスク",
        "与信・ローン審査での人種・地域バイアス",
        "HR・採用にAIを使用する場合は特に公平性審査が必須",
    ]),
], "G")

create_content_slide(prs, "G-4：個人情報・データプライバシー", [
    ("AIと個人情報保護の重要性", [
        "AIシステムは大量の個人データを処理するため、適切な保護が法的・倫理的義務",
        "違反時のリスク：法的罰則・信頼失墜・ブランド毀損",
    ]),
    ("日本の関連法規制", [
        "個人情報保護法（PIPA）：個人情報の取得・利用・提供・管理の規制",
        "改正個人情報保護法（2022年）：プロファイリング規制・漏洩報告義務の強化",
        "医療・金融分野：特定分野の追加規制あり",
    ]),
    ("AIシステムにおけるデータプライバシー設計", [
        "データ最小化原則：必要最小限のデータのみ収集・利用",
        "目的限定原則：収集目的以外への使用禁止",
        "プライバシーバイデザイン：設計段階からプライバシーを組み込む",
        "匿名化・仮名化：学習・分析データの識別子除去",
        "RAGナレッジベースへの個人情報混入防止：アクセス制御・マスキング",
    ]),
    ("クラウドAIサービス利用時の注意点", [
        "データ処理・保管場所の確認（国外サーバー規制）",
        "サービス利用規約でのデータ学習利用の有無確認",
        "エンタープライズプランでの入力データ除外オプション利用",
    ]),
], "G")

create_content_slide(prs, "G-5：AIの透明性・説明可能性（XAI）", [
    ("XAI（Explainable AI）とは", [
        "AIの判断プロセスを人間が理解・解釈できる形で説明する技術・手法",
        "「ブラックボックス問題」への対処：なぜその結論に至ったかを説明",
    ]),
    ("XAIが重要なビジネスシーン", [
        "与信審査・ローン判定：「なぜ否決か」を顧客・規制当局に説明する義務",
        "医療診断支援：医師がAI推奨の根拠を理解した上で最終判断",
        "採用選考：候補者への説明責任・差別禁止法への対応",
        "法的手続き・コンプライアンス：AIの判断が法的に問われるケース",
    ]),
    ("主要なXAI手法", [
        "LIME（Local Interpretable Model-agnostic Explanations）：個別予測の局所的説明",
        "SHAP（SHapley Additive exPlanations）：各特徴量の寄与度を定量化",
        "Attention可視化：LLMがどの部分に注目したかを可視化",
        "RAG引用：回答の根拠文書・箇所を明示（自然な説明可能性）",
    ]),
    ("LLM時代のXAI実践", [
        "RAGシステムでの引用提示が最も実用的なXAI手法",
        "Chain-of-Thought出力で推論ステップを可視化",
        "Confidence Score（確信度スコア）の提示",
    ]),
], "G")

create_content_slide(prs, "G-6：EU AI Act・国内AI規制動向", [
    ("EU AI Act（EU人工知能規制法）", [
        "2024年施行。世界初の包括的AI規制法",
        "リスクベースアプローチ：AIシステムをリスクレベルで分類・規制",
    ]),
    ("EU AI Actのリスク分類", [
        "禁止（Unacceptable Risk）：社会信用スコア、サブリミナル操作、人物の大規模生体認証",
        "高リスク（High Risk）：採用・信用・医療・重要インフラ → 厳格な規制・認証義務",
        "限定リスク（Limited Risk）：チャットボット等 → 透明性義務（AI明示）",
        "最小リスク（Minimal Risk）：スパムフィルタ等 → 自主規制",
    ]),
    ("日本のAI規制動向", [
        "AI事業者ガイドライン（経済産業省・総務省）：安全・透明・公正なAI利活用指針",
        "AI戦略（政府）：生成AI活用推進と安全性確保の両立",
        "業界別ガイドライン：金融庁・厚生労働省等の分野別指針",
        "個人情報保護法との整合：生成AI利用時の個人データ取扱い",
    ]),
    ("企業の対応ポイント", [
        "使用しているAIシステムのリスク分類を確認",
        "高リスク用途では文書化・監査・人間監督の体制整備",
        "EU市場向けサービスはEU AI Act対応が必須",
        "国内ガイドラインに沿ったAIポリシーの策定・公表",
    ]),
], "G")

# ============================================================
# SECTION H: 導入・運用・継続改善
# ============================================================
create_section_divider(prs, "H", "導入・運用・継続改善",
    "パイロット・ステークホルダー・KPIモニタリング・メンテナンス")

create_content_slide(prs, "H-1：パイロット導入の設計", [
    ("パイロット導入とは", [
        "本格展開の前に限定範囲でAIシステムを実際の業務環境でテスト導入する手法",
        "リスクを限定しながら実用性・ROIを検証し、改善してから拡大展開",
    ]),
    ("パイロット設計の原則", [
        "スコープ限定：特定部門・チーム・地域・ユースケースに絞る",
        "代表性確保：本番環境の典型的なケース・ユーザーを含める",
        "期間設定：4〜12週間（長すぎず短すぎず）",
        "測定計画：KPIベースラインと測定方法を事前設定",
        "フィードバック機構：ユーザーからの改善意見収集の仕組み",
    ]),
    ("パイロット成功の判断基準（Go/No-Go）", [
        "精度・品質指標が閾値を超えているか",
        "ユーザー受容性（利用率・満足度）が許容範囲か",
        "ROI見込みが投資を正当化できるか",
        "セキュリティ・コンプライアンス問題が発生していないか",
        "運用・保守体制が確立できているか",
    ]),
    ("パイロットから本番展開への移行", [
        "パイロット結果の文書化・共有（成功・失敗両方のナレッジ）",
        "スケールアップ計画の策定（展開範囲・タイムライン・予算）",
        "本番インフラへの移行・負荷テスト",
        "ロールアウトコミュニケーション計画の実行",
    ]),
], "H")

create_content_slide(prs, "H-2：ステークホルダーマネジメント", [
    ("ステークホルダーマネジメントとは", [
        "AI導入プロジェクトに関わる・影響を受ける全ての関係者を特定し、効果的に関与・管理する活動",
    ]),
    ("主要ステークホルダーの特定と分類", [
        "経営層（スポンサー）：戦略的支援・予算承認・変革のメッセージ発信",
        "業務部門長（推進者）：業務要件定義・現場推進・変化管理",
        "IT・システム部門：技術実装・セキュリティ・インフラ管理",
        "現場ユーザー：実際の利用者・フィードバック提供者",
        "法務・コンプライアンス：規制対応・リスク審査",
        "外部ベンダー・パートナー：技術サポート・ライセンス管理",
    ]),
    ("ステークホルダー関与戦略（影響力×関心度マトリクス）", [
        "高影響力×高関心度 → 密接な協力・定期的な関与（Key Stakeholders）",
        "高影響力×低関心度 → 満足状態の維持・重要情報の提供（Keep Satisfied）",
        "低影響力×高関心度 → 十分な情報提供・意見収集（Keep Informed）",
        "低影響力×低関心度 → 最低限の情報共有（Monitor）",
    ]),
    ("コミュニケーション計画", [
        "誰に・何を・いつ・どの媒体で伝えるかを計画化",
        "定期報告（週次進捗・月次KPIレポート）の仕組み化",
        "課題・リスクのエスカレーションルートの明確化",
    ]),
], "H")

create_content_slide(prs, "H-3：変化管理（チェンジマネジメント）", [
    ("AI導入における変化管理の3要素", [
        "①意識変革：「AIとの協働」という新しい働き方への理解促進",
        "②スキル変革：AIツールを活用するための能力開発・トレーニング",
        "③行動変革：新しいプロセス・ツールを日常業務で実際に使うようになる",
    ]),
    ("抵抗管理：なぜ人はAI導入に抵抗するか", [
        "雇用への不安：「自分の仕事がなくなる」という恐怖",
        "能力不安：「使いこなせるか」というスキルへの不安",
        "プライバシー懸念：「自分の業務データが監視される」という警戒",
        "慣性：「今の方法で十分うまくいっている」という惰性",
    ]),
    ("変化管理の実践策", [
        "早期巻き込み：設計段階から現場ユーザーを参加させる",
        "クイックウィン：小さな成功体験を作り、メリットを体感させる",
        "ピアチャンピオン：同じ現場の「AI推進者」が同僚を支援",
        "継続トレーニング：実践的・繰り返しの学習機会提供",
        "心理的安全性：失敗を許容し試行錯誤を奨励する文化",
    ]),
    ("評価制度との連携", [
        "AI活用を業績評価・昇格要件に組み込む",
        "AI推進への貢献を表彰・インセンティブとして認める",
        "デジタルスキルを採用・育成基準に追加",
    ]),
], "H")

create_content_slide(prs, "H-4：KPIモニタリングと改善サイクル", [
    ("継続的モニタリングの必要性", [
        "AIシステムはデプロイ後も環境変化（データドリフト・業務変化）で性能が変化する",
        "定期的な測定・評価・改善サイクルが持続的な価値創出に不可欠",
    ]),
    ("モニタリングすべき指標", [
        "技術指標：精度・F1スコア・応答時間・エラー率・API コスト",
        "業務指標：業務処理時間・エラー件数・対応件数・コスト削減額",
        "ユーザー指標：利用率・ユーザー満足度（NPS）・フィードバック数",
        "リスク指標：ハルシネーション発生率・セキュリティインシデント・コンプライアンス違反",
    ]),
    ("データドリフトとモデル劣化", [
        "データドリフト：入力データの分布が学習時と変化する現象",
        "コンセプトドリフト：業務ルール・市場環境の変化でモデルの前提が崩れる",
        "検出方法：入力分布の統計的モニタリング・精度の定期評価",
        "対応：RAGナレッジベースの定期更新・プロンプト改善・再学習",
    ]),
    ("改善サイクル（PDCAサイクル）", [
        "Plan：KPI目標設定・改善計画策定",
        "Do：改善施策の実施（プロンプト改善・データ更新・設定変更）",
        "Check：KPI計測・ユーザーフィードバック収集",
        "Act：評価に基づく次の改善計画へ",
    ]),
], "H")

create_content_slide(prs, "H-5：AIエージェントのメンテナンス", [
    ("AIエージェントのメンテナンスの特徴", [
        "従来のソフトウェアと異なり、AI特有のメンテナンス要素が存在する",
        "継続的な改善が競争優位性の源泉",
    ]),
    ("プロンプトのメンテナンス", [
        "業務変化・新要件に合わせたプロンプトの更新",
        "ユーザーフィードバックに基づく改善・最適化",
        "プロンプトのバージョン管理（変更履歴・ロールバック可能に）",
        "A/Bテスト：複数バージョンの比較・最良版の選定",
    ]),
    ("RAGナレッジベースのメンテナンス", [
        "定期的な文書の追加・更新・削除（鮮度管理）",
        "不正確・古い情報の除去（データクリーニング）",
        "新規業務ルール・製品情報の即時追加",
        "チャンクの再最適化・インデックスの再構築",
    ]),
    ("LLMモデルのアップデート対応", [
        "ベースモデルのバージョンアップ時の動作検証",
        "APIの仕様変更への追従",
        "新機能（マルチモーダル等）の活用検討",
    ]),
    ("インシデント対応とエスカレーション", [
        "障害・誤動作発生時の検知→報告→対応フローの整備",
        "ロールバック計画：旧バージョンへの切り戻し手順",
        "事後レビュー（ポストモーテム）：原因分析と再発防止",
    ]),
], "H")

# ============================================================
# SUMMARY / EXAM TIPS SLIDE
# ============================================================
create_content_slide(prs, "試験対策：重要ポイントまとめ", [
    ("試験の特徴と対策方針", [
        "4択ケーススタディ形式：実際のビジネス場面でどう判断・行動するかを問う",
        "知識の暗記ではなく「なぜそうするか」という理解が問われる",
        "全8領域（A〜H）を横断した統合的理解が重要",
    ]),
    ("頻出テーマと要注意ポイント", [
        "ECRSの順序：Eliminate→Combine→Rearrange→Simplify（この順で検討）",
        "RAGの目的3点：幻覚抑制・最新性確保・根拠提示（セットで記憶）",
        "5Dモデルの順序：Discovery→Definition→Design→Development→Deployment",
        "AIデリゲーション判断：频度・ルール化・可逆性・説明責任の4基準",
        "組織モデル：ハイブリッド型が多くのケースで最適解として問われる",
        "EU AI Act：リスク分類（禁止・高リスク・限定リスク・最小リスク）",
    ]),
    ("実践的思考の鍛え方", [
        "各ユースケースに対してIPOフレームワークで分解する習慣",
        "「まずECRSで最適化してからAI化」の思考順序を守る",
        "スポンサーシップの重要性：技術より人・組織の問題が失敗原因になる",
        "ハルシネーションリスクへの対策としてRAGが最も実用的な解決策",
    ]),
], "")

# Final slide
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)
add_background(slide, NAVY)
add_rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.12), GOLD)
add_rect(slide, Inches(0), Inches(7.38), Inches(13.33), Inches(0.12), GOLD)

add_text_box(slide, "AIエージェント・ストラテジスト", Inches(1), Inches(2.0), Inches(11), Inches(1.0),
            font_size=34, font_color=GOLD, bold=True, align=PP_ALIGN.CENTER)
add_text_box(slide, "認定試験 合格を目指して", Inches(1), Inches(3.1), Inches(11), Inches(0.7),
            font_size=24, font_color=WHITE, bold=False, align=PP_ALIGN.CENTER)
add_text_box(slide, "一般社団法人 AICX協会", Inches(1), Inches(4.0), Inches(11), Inches(0.5),
            font_size=18, font_color=LIGHT_GOLD, align=PP_ALIGN.CENTER)
add_text_box(slide, "初回試験：2026年7月", Inches(1), Inches(4.6), Inches(11), Inches(0.4),
            font_size=16, font_color=RGBColor(0xB0, 0xBE, 0xC5), align=PP_ALIGN.CENTER)

# Save the presentation
output_path = "/home/user/claude-code-on-the-web-html/ai_agent_strategist_exam.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
