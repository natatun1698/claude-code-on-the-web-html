#!/usr/bin/env python3
"""Microsoft AI Tour Tokyo 基調講演・西脇資哲さん発表(最後の17分)のスクリーンショットPPTXを生成する。"""
import glob
import os
import re
import zipfile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

VIDEO_ID = "5r90rz_HQ38"
VIDEO_URL = f"https://www.youtube.com/live/{VIDEO_ID}"
FRAMES = "frames_up"
OUT = "nishiwaki_last17min.pptx"

NAVY = RGBColor(0x0B, 0x14, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xC9, 0xD6, 0xEF)
ACCENT = RGBColor(0x4C, 0xC2, 0xFF)
GRAY = RGBColor(0x8A, 0x9A, 0xB8)

SW, SH = Inches(13.333), Inches(7.5)

files = sorted(glob.glob(f"{FRAMES}/*.png"))
times = [int(os.path.basename(f)[1:5]) for f in files]


def nearest(t):
    i = min(range(len(times)), key=lambda i: abs(times[i] - t))
    return files[i], times[i]


def mmss(t):
    return f"{t // 60}:{t % 60:02d}"


prs = Presentation()
prs.slide_width = SW
prs.slide_height = SH
blank = prs.slide_layouts[6]


def add_slide():
    s = prs.slides.add_slide(blank)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = NAVY
    return s


def txt(slide, x, y, w, h, text, size, color=WHITE, bold=False, align=PP_ALIGN.LEFT, link=None, font="Yu Gothic UI"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    first = True
    for line in text.split("\n"):
        if not first:
            p = tf.add_paragraph()
            p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
        if link:
            r.hyperlink.address = link
        first = False
    return tb


def grid_slide(title, sub, items, cols=3, note=None):
    """items: list of (target_sec, caption)"""
    s = add_slide()
    txt(s, Inches(0.45), Inches(0.16), Inches(12.4), Inches(0.55), title, 24, WHITE, bold=True)
    if sub:
        txt(s, Inches(0.45), Inches(0.72), Inches(12.4), Inches(0.38), sub, 13, LIGHT)
    rows = (len(items) + cols - 1) // cols
    gy = Inches(1.22)
    gap = Emu(60000)
    avail_h = SH - gy - (Inches(0.95) if note else Inches(0.25))
    cell_h = int(avail_h / rows)
    img_h = cell_h - Inches(0.30)
    cell_w = int(img_h * 16 / 9)
    max_w = int((SW - Inches(0.6) - gap * (cols - 1)) / cols)
    if cell_w > max_w:
        cell_w = max_w
        img_h = int(cell_w * 9 / 16)
        cell_h = img_h + Inches(0.30)
    total_w = cell_w * cols + gap * (cols - 1)
    gx = int((SW - total_w) / 2)
    for i, (t, cap) in enumerate(items):
        f, tt = nearest(t)
        x = gx + (i % cols) * (cell_w + gap)
        y = gy + (i // cols) * cell_h
        s.shapes.add_picture(f, x, y, width=cell_w, height=img_h)
        label = mmss(tt) + (f"  {cap}" if cap else "")
        txt(s, x, y + img_h + Emu(9525), cell_w, Inches(0.26), label, 10, ACCENT,
            link=f"{VIDEO_URL}?t={tt}s")
    if note:
        txt(s, Inches(0.45), SH - Inches(0.88), Inches(12.4), Inches(0.75), note, 12, LIGHT)
    return s


# ---------- 1. 表紙 ----------
s = add_slide()
f, t = nearest(5182)
s.shapes.add_picture(f, Inches(7.1), Inches(1.5), width=Inches(5.6))
txt(s, Inches(0.55), Inches(1.15), Inches(6.6), Inches(0.5),
    "Microsoft AI Tour Tokyo 基調講演", 18, ACCENT, bold=True)
txt(s, Inches(0.55), Inches(1.75), Inches(6.4), Inches(1.9),
    "西脇資哲さんの発表\n(最後の17分)", 40, WHITE, bold=True)
txt(s, Inches(0.55), Inches(3.9), Inches(6.4), Inches(1.6),
    "「日本の AI フロンティアを切り開く」\n2026年3月24日 東京ビッグサイト\n動画 1:25:48 〜 1:42:48 のダイジェスト", 15, LIGHT)
txt(s, Inches(0.55), Inches(5.6), Inches(6.6), Inches(0.4),
    "▶ 元動画: " + VIDEO_URL, 12, ACCENT, link=VIDEO_URL)
txt(s, Inches(7.1), Inches(4.75), Inches(5.6), Inches(0.35),
    f"{mmss(t)} 登壇中の西脇資哲氏(日本マイクロソフト エバンジェリスト)", 10, GRAY)

# ---------- 2. この資料について ----------
s = add_slide()
txt(s, Inches(0.45), Inches(0.3), Inches(12.4), Inches(0.6), "この資料について", 26, WHITE, bold=True)
body = (
    "・本資料はYouTubeアーカイブの最後の17分間(1:25:48〜1:42:48、西脇資哲氏の発表パート)を、\n"
    "   実際の映像フレーム(約10秒間隔)で振り返るダイジェストです。\n"
    "\n"
    "・画像はYouTubeが公開しているシークプレビュー画像(ストーリーボード、原寸160×90px)を拡大したものです。\n"
    "   実行環境からの動画本体のダウンロードがYouTube側の制限(データセンターIPのボット判定)で行えなかったため、\n"
    "   解像度が低く細かい文字は判読できません。\n"
    "\n"
    "・各画像の下のタイムスタンプはリンクになっており、クリックするとYouTubeの該当シーンが開きます。\n"
    "   高解像度で確認したい場面はリンク先をご覧ください。\n"
    "\n"
    "・各セクションの説明文は、映像フレームから読み取れる内容と公開イベントレポートに基づく要約であり、\n"
    "   一部推定を含みます。"
)
txt(s, Inches(0.7), Inches(1.3), Inches(12.0), Inches(5.5), body, 15, LIGHT)

# ---------- 3. タイムライン ----------
s = add_slide()
txt(s, Inches(0.45), Inches(0.25), Inches(12.4), Inches(0.6), "最後の17分の流れ", 26, WHITE, bold=True)
timeline = [
    ("1:25:48", "対談パートの締め 〜 西脇資哲氏 登壇", 5148),
    ("1:26:13", "タイトル「西脇資哲 エバンジェリスト」/ 課題提起:書類の山と時間", 5173),
    ("1:28:23", "デモ①  Excel × Copilot ― 売上データの集計・分析・グラフ化", 5303),
    ("1:31:21", "デモ②  PowerPoint × Copilot ―「2025年度 決算説明会」資料の自動生成", 5481),
    ("1:34:30", "デモ③  リサーチ&営業準備 ― シャープ(株)の企業調査から提案資料まで", 5670),
    ("1:38:59", "クロージング ― Microsoftのミッションと日本へのメッセージ", 5939),
    ("1:40:18", "エンディング映像(日本のものづくりの現場) 〜 終了", 6018),
]
y = Inches(1.15)
for tm, desc, sec in timeline:
    txt(s, Inches(0.7), y, Inches(1.5), Inches(0.4), tm, 15, ACCENT, bold=True,
        link=f"{VIDEO_URL}?t={sec}s")
    txt(s, Inches(2.3), y, Inches(10.4), Inches(0.4), desc, 15, WHITE)
    y += Inches(0.52)
strip = [5173, 5352, 5600, 5849, 5988, 6068]
x = Inches(0.7)
for tsec in strip:
    f, tt = nearest(tsec)
    s.shapes.add_picture(f, x, Inches(5.15), width=Inches(1.95))
    txt(s, x, Inches(6.28), Inches(1.95), Inches(0.25), mmss(tt), 9, ACCENT,
        link=f"{VIDEO_URL}?t={tt}s")
    x += Inches(2.05)

# ---------- 登壇 ----------
grid_slide(
    "登壇 ― 西脇資哲氏",
    "1:25:48〜1:28:20   日本マイクロソフト 業務執行役員 エバンジェリスト。基調講演の締めくくりとして登壇し、「AIで日々の仕事がどう変わるか」を実演中心に紹介",
    [(5143, "前パートからの転換"), (5163, "登壇"), (5173, "タイトル表示「西脇 資哲」"),
     (5222, "トーク"), (5262, "トーク"), (5292, "課題提起: 時計と書類の山")],
    cols=3,
    note="冒頭では「時計と積み上がる書類」のビジュアルで、資料作成・情報収集に費やされる膨大な時間という職場の課題を提起。ここからCopilotによる時短デモが始まります。",
)

# ---------- デモ1 Excel (2枚) ----------
grid_slide(
    "デモ①  Excel × Copilot ― データ集計・分析 (1/2)",
    "1:28:23〜1:29:50   大きな売上データをCopilotとの対話だけで整形・分析",
    [(5303, "売上データを開く"), (5323, "Copilotに指示"), (5343, "表の整形"),
     (5362, "集計処理"), (5382, "計算結果の反映"), (5401, "データのハイライト")],
    cols=3,
    note="Excel上のCopilotに自然言語で指示し、大きな売上データの整形・集計・条件付き強調を対話だけで実行していきます。",
)
grid_slide(
    "デモ①  Excel × Copilot ― データ集計・分析 (2/2)",
    "1:29:50〜1:31:10   分析からグラフ・ピボットまで一気に生成",
    [(5411, "分析の実行"), (5421, "「ブックを編集しましょう」"), (5431, "Copilotが処理中"),
     (5441, "集計ビュー"), (5451, "グラフ生成"), (5461, "チャート/ピボット完成")],
    cols=3,
    note="手作業なら数時間かかる分析・可視化が数分で完了する様子をライブでデモ。円グラフや棒グラフを含むダッシュボード風の出力まで一気に到達します。",
)

# ---------- デモ2 PowerPoint (2枚) ----------
grid_slide(
    "デモ②  PowerPoint × Copilot ― 決算説明会資料 (1/2)",
    "1:31:11〜1:33:10   Excelの分析結果からプレゼン資料をCopilotが自動作成",
    [(5471, "再び「時間」の課題へ"), (5491, "PowerPoint × Copilot"), (5511, "生成の指示"),
     (5531, "構成案の生成"), (5551, "ビジュアル素材の提示"), (5571, "スライド生成中")],
    cols=3,
    note="「決算説明会の資料を作って」という指示から、Copilotが構成案・ビジュアルを含むスライド一式のドラフトを生成していきます。",
)
grid_slide(
    "デモ②  PowerPoint × Copilot ― 決算説明会資料 (2/2)",
    "1:33:10〜1:34:20   「2025年度 決算説明会」資料が完成",
    [(5591, "デッキの仕上がり"), (5601, "表紙「2025年度 決算説明会」"), (5611, "登壇者の解説"),
     (5621, "内容スライド"), (5641, "KPI強調(14.3%)のビジュアル化"), (5651, "資料の確認")],
    cols=3,
    note="文字だらけのスライドも「視覚的に表現して」と指示するだけで、数値を大きく見せるモダンなレイアウトに変換。経営層向け資料が短時間で仕上がることを示しました。",
)

# ---------- デモ3 リサーチ (2枚) ----------
grid_slide(
    "デモ③  リサーチ & 営業準備 (1/2)",
    "1:34:30〜1:37:00   Copilotのリサーチ機能で訪問前の営業準備を自動化",
    [(5670, "デモ転換"), (5700, "リサーチの指示"), (5720, "調査レポート生成中"),
     (5750, "レポート本文"), (5780, "詳細の深掘り"), (5810, "出典付きの調査結果")],
    cols=3,
    note="訪問先企業の情報・最新動向をCopilot(リサーチ機能)が出典付きで深掘り調査。長文レポートが自動で組み上がっていきます。",
)
grid_slide(
    "デモ③  リサーチ & 営業準備 (2/2)",
    "1:37:00〜1:38:50   シャープ株式会社の企業調査から提案アプローチまで",
    [(5830, "調査のまとめ"), (5849, "「シャープ株式会社 企業情報」"), (5869, "関連情報の整理"),
     (5889, "「なぜ今、この人事異動か」"), (5909, "社長への提案と注意点"), (5929, "提案アプローチ戦略")],
    cols=3,
    note="シャープ株式会社を例に、企業情報・人事異動の分析・つながり(LinkedIn)情報・提案アプローチまで自動整理。営業担当の事前準備が一変することを示しました。",
)

# ---------- クロージング ----------
grid_slide(
    "クロージング ― Microsoftのミッションと日本へのメッセージ",
    "1:38:59〜1:40:10",
    [(5939, "ステージ全景"), (5949, "締めのトーク"), (5968, "会場へのメッセージ"),
     (5988, "ミッションスライド"), (5998, "「あなたのCopilotとして」"), (6008, "発表の結び")],
    cols=3,
    note="「地球上のすべての個人とすべての組織が、より多くのことを達成できるようにする」というMicrosoftのミッションを掲げ、日本マイクロソフトが皆さんの“Copilot”として伴走するというメッセージで発表を締めくくりました。",
)

# ---------- エンディング映像 ----------
grid_slide(
    "エンディング映像 ― 日本のものづくりの現場へ",
    "1:40:18〜1:42:48   基調講演の幕引きとなるブランドムービー 〜 Microsoftロゴで終了",
    [(6018, "映像スタート"), (6038, "工場の現場"), (6058, "製造ライン"),
     (6078, "現場のディテール"), (6118, "働く人の姿"), (6157, "Microsoftロゴで終了")],
    cols=3,
    note="日本の製造業の現場を映したエンディングムービーが流れ、Microsoftロゴとともに基調講演は終了。AIが日本の現場の力を引き出すという基調講演全体のテーマを象徴する締めくくりでした。",
)

# ---------- 付録モンタージュ ----------
chunks = [files[i:i + 35] for i in range(0, len(files), 35)]
for ci, chunk in enumerate(chunks):
    s = add_slide()
    txt(s, Inches(0.45), Inches(0.18), Inches(12.4), Inches(0.5),
        f"付録: 全フレーム一覧 ({ci + 1}/{len(chunks)})  ―  約10秒間隔のストーリーボード", 20, WHITE, bold=True)
    cols = 7
    cell_w = Inches(1.78)
    img_h = Inches(1.0)
    gx, gy = Inches(0.45), Inches(0.85)
    for i, f in enumerate(chunk):
        t = int(os.path.basename(f)[1:5])
        x = gx + (i % cols) * (cell_w + Inches(0.06))
        y = gy + (i // cols) * (img_h + Inches(0.30))
        s.shapes.add_picture(f, x, y, width=cell_w, height=img_h)
        txt(s, x, y + img_h, cell_w, Inches(0.22), mmss(t), 8, ACCENT,
            link=f"{VIDEO_URL}?t={t}s")

prs.save(OUT)

# テーマのハイパーリンク色を暗背景でも読めるアクセント色に差し替える
tmp = OUT + ".tmp"
with zipfile.ZipFile(OUT) as zin, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
    for item in zin.namelist():
        data = zin.read(item)
        if item.endswith("theme1.xml"):
            xml = data.decode("utf-8")
            xml = re.sub(r'(<a:hlink>\s*<a:srgbClr val=")[0-9A-Fa-f]{6}("/>)', r"\g<1>4CC2FF\g<2>", xml)
            xml = re.sub(r'(<a:folHlink>\s*<a:srgbClr val=")[0-9A-Fa-f]{6}("/>)', r"\g<1>9AD8FF\g<2>", xml)
            data = xml.encode("utf-8")
        zout.writestr(item, data)
os.replace(tmp, OUT)
print("saved", OUT, os.path.getsize(OUT), "bytes,", len(prs.slides._sldIdLst), "slides")
